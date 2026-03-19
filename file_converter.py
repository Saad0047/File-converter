"""
File Converter - Python Desktop App
------------------------------------
Converts between popular file formats using a GUI file picker.

Supported conversions:
  DOCX  → PDF, TXT, HTML, Markdown
  PDF   → TXT, DOCX (basic), HTML
  TXT   → PDF, DOCX, HTML, Markdown
  HTML  → TXT, PDF, Markdown, DOCX
  MD    → HTML, PDF, TXT, DOCX
  CSV   → JSON, XLSX, TXT, HTML
  JSON  → CSV, TXT, XLSX
  XLSX  → CSV, JSON, TXT
  PNG/JPG/JPEG/BMP/TIFF/WEBP → Any other image format, PDF
  PPTX  → PDF, TXT
  XML   → JSON, TXT

Install dependencies:
  pip install python-docx reportlab markdown beautifulsoup4 lxml
              openpyxl pillow pypdf2 python-pptx
"""

import tkinter as tk
from tkinter import filedialog, messagebox, ttk
import os
import json
import csv
import threading


# ─── Dependency check ───────────────────────────────────────────────────────

def check_deps():
    missing = []
    for pkg, imp in [
        ("python-docx",    "docx"),
        ("reportlab",      "reportlab"),
        ("markdown",       "markdown"),
        ("beautifulsoup4", "bs4"),
        ("openpyxl",       "openpyxl"),
        ("Pillow",         "PIL"),
        ("pypdf2",         "PyPDF2"),
        ("python-pptx",    "pptx"),
        ("lxml",           "lxml"),
    ]:
        try:
            __import__(imp)
        except ImportError:
            missing.append(pkg)
    return missing


# ─── Conversion logic ────────────────────────────────────────────────────────

def convert_file(src, target_fmt):
    """Main dispatcher. Returns (output_path, None) or (None, error_string)."""
    ext = os.path.splitext(src)[1].lower().lstrip(".")
    base = os.path.splitext(src)[0]
    dst  = f"{base}_converted.{target_fmt}"

    try:
        key = (ext, target_fmt)

        # ── DOCX conversions ──────────────────────────────────────────────
        if key == ("docx", "txt"):
            import docx
            doc = docx.Document(src)
            text = "\n".join(p.text for p in doc.paragraphs)
            _write_text(dst, text)

        elif key == ("docx", "html"):
            import docx
            doc = docx.Document(src)
            paragraphs = [f"<p>{p.text}</p>" for p in doc.paragraphs if p.text.strip()]
            html = _html_page("\n".join(paragraphs))
            _write_text(dst, html)

        elif key == ("docx", "md"):
            import docx
            doc = docx.Document(src)
            lines = []
            for p in doc.paragraphs:
                if p.style.name.startswith("Heading 1"):
                    lines.append(f"# {p.text}")
                elif p.style.name.startswith("Heading 2"):
                    lines.append(f"## {p.text}")
                elif p.style.name.startswith("Heading 3"):
                    lines.append(f"### {p.text}")
                else:
                    lines.append(p.text)
            _write_text(dst, "\n".join(lines))

        elif key == ("docx", "pdf"):
            import docx
            from reportlab.lib.pagesizes import A4
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet
            doc = docx.Document(src)
            pdf = SimpleDocTemplate(dst, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            for p in doc.paragraphs:
                if p.text.strip():
                    story.append(Paragraph(p.text, styles["Normal"]))
                    story.append(Spacer(1, 6))
            pdf.build(story)

        # ── PDF conversions ───────────────────────────────────────────────
        elif key == ("pdf", "txt"):
            import PyPDF2
            text = ""
            with open(src, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() or ""
            _write_text(dst, text)

        elif key == ("pdf", "html"):
            import PyPDF2
            text = ""
            with open(src, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() or ""
            paragraphs = "\n".join(f"<p>{line}</p>" for line in text.splitlines() if line.strip())
            _write_text(dst, _html_page(paragraphs))

        elif key == ("pdf", "docx"):
            import PyPDF2
            import docx
            text = ""
            with open(src, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() or ""
            doc = docx.Document()
            for line in text.splitlines():
                doc.add_paragraph(line)
            doc.save(dst)

        # ── TXT conversions ───────────────────────────────────────────────
        elif key == ("txt", "pdf"):
            text = _read_text(src)
            _text_to_pdf(text, dst)

        elif key == ("txt", "docx"):
            import docx
            doc = docx.Document()
            for line in _read_text(src).splitlines():
                doc.add_paragraph(line)
            doc.save(dst)

        elif key == ("txt", "html"):
            text = _read_text(src)
            body = "\n".join(f"<p>{line or '&nbsp;'}</p>" for line in text.splitlines())
            _write_text(dst, _html_page(body))

        elif key == ("txt", "md"):
            import shutil
            shutil.copy(src, dst)   # Plain text is already valid Markdown

        # ── Markdown conversions ──────────────────────────────────────────
        elif key == ("md", "html"):
            import markdown
            text = _read_text(src)
            html = markdown.markdown(text, extensions=["tables", "fenced_code"])
            _write_text(dst, _html_page(html))

        elif key == ("md", "txt"):
            from bs4 import BeautifulSoup
            import markdown
            html = markdown.markdown(_read_text(src))
            soup = BeautifulSoup(html, "lxml")
            _write_text(dst, soup.get_text())

        elif key == ("md", "pdf"):
            from bs4 import BeautifulSoup
            import markdown
            html = markdown.markdown(_read_text(src))
            soup = BeautifulSoup(html, "lxml")
            _text_to_pdf(soup.get_text(), dst)

        elif key == ("md", "docx"):
            from bs4 import BeautifulSoup
            import markdown, docx
            html = markdown.markdown(_read_text(src))
            soup = BeautifulSoup(html, "lxml")
            doc = docx.Document()
            for tag in soup.find_all(["h1","h2","h3","p","li"]):
                if tag.name == "h1":
                    doc.add_heading(tag.get_text(), level=1)
                elif tag.name == "h2":
                    doc.add_heading(tag.get_text(), level=2)
                elif tag.name == "h3":
                    doc.add_heading(tag.get_text(), level=3)
                else:
                    doc.add_paragraph(tag.get_text())
            doc.save(dst)

        # ── HTML conversions ──────────────────────────────────────────────
        elif key == ("html", "txt"):
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(_read_text(src), "lxml")
            _write_text(dst, soup.get_text())

        elif key == ("html", "pdf"):
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(_read_text(src), "lxml")
            _text_to_pdf(soup.get_text(), dst)

        elif key == ("html", "md"):
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(_read_text(src), "lxml")
            lines = []
            for tag in soup.find_all(["h1","h2","h3","h4","p","li","br"]):
                t = tag.get_text().strip()
                if not t:
                    continue
                if tag.name == "h1":   lines.append(f"# {t}")
                elif tag.name == "h2": lines.append(f"## {t}")
                elif tag.name == "h3": lines.append(f"### {t}")
                elif tag.name == "h4": lines.append(f"#### {t}")
                elif tag.name == "li": lines.append(f"- {t}")
                else:                  lines.append(t)
            _write_text(dst, "\n".join(lines))

        elif key == ("html", "docx"):
            from bs4 import BeautifulSoup
            import docx
            soup = BeautifulSoup(_read_text(src), "lxml")
            doc = docx.Document()
            for tag in soup.find_all(["h1","h2","h3","p","li"]):
                t = tag.get_text().strip()
                if not t:
                    continue
                if tag.name in ("h1","h2","h3"):
                    doc.add_heading(t, level=int(tag.name[1]))
                else:
                    doc.add_paragraph(t)
            doc.save(dst)

        # ── CSV conversions ───────────────────────────────────────────────
        elif key == ("csv", "json"):
            with open(src, newline="", encoding="utf-8") as f:
                rows = list(csv.DictReader(f))
            _write_text(dst, json.dumps(rows, indent=2, ensure_ascii=False))

        elif key == ("csv", "xlsx"):
            import openpyxl
            wb = openpyxl.Workbook()
            ws = wb.active
            with open(src, newline="", encoding="utf-8") as f:
                for row in csv.reader(f):
                    ws.append(row)
            wb.save(dst)

        elif key == ("csv", "txt"):
            import shutil
            shutil.copy(src, dst)

        elif key == ("csv", "html"):
            with open(src, newline="", encoding="utf-8") as f:
                reader = csv.reader(f)
                rows = list(reader)
            if not rows:
                raise ValueError("Empty CSV")
            header = "".join(f"<th>{c}</th>" for c in rows[0])
            body   = "".join(
                "<tr>" + "".join(f"<td>{cell}</td>" for cell in row) + "</tr>"
                for row in rows[1:]
            )
            table = f"<table border='1'><thead><tr>{header}</tr></thead><tbody>{body}</tbody></table>"
            _write_text(dst, _html_page(table))

        # ── JSON conversions ──────────────────────────────────────────────
        elif key == ("json", "csv"):
            with open(src, encoding="utf-8") as f:
                data = json.load(f)
            if not isinstance(data, list):
                data = [data]
            if not data:
                raise ValueError("Empty JSON array")
            keys = list(data[0].keys())
            with open(dst, "w", newline="", encoding="utf-8") as f:
                writer = csv.DictWriter(f, fieldnames=keys, extrasaction="ignore")
                writer.writeheader()
                writer.writerows(data)

        elif key == ("json", "txt"):
            _write_text(dst, json.dumps(json.load(open(src)), indent=2))

        elif key == ("json", "xlsx"):
            import openpyxl
            with open(src, encoding="utf-8") as f:
                data = json.load(f)
            if not isinstance(data, list):
                data = [data]
            wb = openpyxl.Workbook()
            ws = wb.active
            if data:
                ws.append(list(data[0].keys()))
                for row in data:
                    ws.append(list(row.values()))
            wb.save(dst)

        # ── XLSX conversions ──────────────────────────────────────────────
        elif key == ("xlsx", "csv"):
            import openpyxl
            wb = openpyxl.load_workbook(src, data_only=True)
            ws = wb.active
            with open(dst, "w", newline="", encoding="utf-8") as f:
                writer = csv.writer(f)
                for row in ws.iter_rows(values_only=True):
                    writer.writerow(row)

        elif key == ("xlsx", "json"):
            import openpyxl
            wb = openpyxl.load_workbook(src, data_only=True)
            ws = wb.active
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                raise ValueError("Empty spreadsheet")
            headers = [str(c) for c in rows[0]]
            result = [dict(zip(headers, row)) for row in rows[1:]]
            _write_text(dst, json.dumps(result, indent=2, default=str))

        elif key == ("xlsx", "txt"):
            import openpyxl
            wb = openpyxl.load_workbook(src, data_only=True)
            ws = wb.active
            lines = ["\t".join(str(c) if c is not None else "" for c in row)
                     for row in ws.iter_rows(values_only=True)]
            _write_text(dst, "\n".join(lines))

        # ── Image conversions ─────────────────────────────────────────────
        elif ext in ("png","jpg","jpeg","bmp","tiff","webp","gif") and \
             target_fmt in ("png","jpg","jpeg","bmp","tiff","webp","pdf"):
            from PIL import Image
            img = Image.open(src)
            if target_fmt == "pdf":
                rgb = img.convert("RGB")
                rgb.save(dst, "PDF", resolution=150)
            else:
                save_fmt = "JPEG" if target_fmt in ("jpg","jpeg") else target_fmt.upper()
                save_img = img.convert("RGB") if save_fmt == "JPEG" else img
                save_img.save(dst, save_fmt)

        # ── PPTX conversions ──────────────────────────────────────────────
        elif key == ("pptx", "txt"):
            from pptx import Presentation
            prs = Presentation(src)
            lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            lines.append(para.text)
                lines.append("")
            _write_text(dst, "\n".join(lines))

        elif key == ("pptx", "pdf"):
            from pptx import Presentation
            from reportlab.lib.pagesizes import A4
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet
            prs = Presentation(src)
            lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                lines.append(para.text)
            _text_to_pdf("\n".join(lines), dst)

        # ── XML conversions ───────────────────────────────────────────────
        elif key == ("xml", "json"):
            import xml.etree.ElementTree as ET
            def elem_to_dict(elem):
                d = {elem.tag: {} if elem.attrib or list(elem) else elem.text}
                if elem.attrib:
                    d[elem.tag].update({"@" + k: v for k, v in elem.attrib.items()})
                for child in elem:
                    child_d = elem_to_dict(child)
                    if child.tag in d[elem.tag]:
                        existing = d[elem.tag][child.tag]
                        if not isinstance(existing, list):
                            d[elem.tag][child.tag] = [existing]
                        d[elem.tag][child.tag].append(child_d[child.tag])
                    else:
                        d[elem.tag].update(child_d)
                return d
            tree = ET.parse(src)
            _write_text(dst, json.dumps(elem_to_dict(tree.getroot()), indent=2))

        elif key == ("xml", "txt"):
            import xml.etree.ElementTree as ET
            tree = ET.parse(src)
            texts = [elem.text for elem in tree.iter() if elem.text and elem.text.strip()]
            _write_text(dst, "\n".join(texts))

        else:
            return None, f"Conversion from .{ext} to .{target_fmt} is not supported."

        return dst, None

    except Exception as e:
        return None, str(e)


# ─── Helpers ─────────────────────────────────────────────────────────────────

def _read_text(path):
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()

def _write_text(path, content):
    with open(path, "w", encoding="utf-8") as f:
        f.write(content)

def _html_page(body):
    return f"""<!DOCTYPE html>
<html lang="en">
<head><meta charset="UTF-8"><title>Converted</title>
<style>body{{font-family:sans-serif;max-width:800px;margin:2rem auto;line-height:1.6}}</style>
</head>
<body>{body}</body>
</html>"""

def _text_to_pdf(text, dst):
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    doc = SimpleDocTemplate(dst, pagesize=A4)
    styles = getSampleStyleSheet()
    story = []
    for line in text.splitlines():
        story.append(Paragraph(line or "&nbsp;", styles["Normal"]))
        story.append(Spacer(1, 4))
    doc.build(story)


# ─── Format map (what each input can convert to) ──────────────────────────────

FORMAT_TARGETS = {
    "docx":  ["pdf", "txt", "html", "md"],
    "pdf":   ["txt", "html", "docx"],
    "txt":   ["pdf", "docx", "html", "md"],
    "md":    ["html", "pdf", "txt", "docx"],
    "html":  ["txt", "pdf", "md", "docx"],
    "csv":   ["json", "xlsx", "txt", "html"],
    "json":  ["csv", "txt", "xlsx"],
    "xlsx":  ["csv", "json", "txt"],
    "png":   ["jpg", "bmp", "tiff", "webp", "pdf"],
    "jpg":   ["png", "bmp", "tiff", "webp", "pdf"],
    "jpeg":  ["png", "bmp", "tiff", "webp", "pdf"],
    "bmp":   ["png", "jpg", "tiff", "webp", "pdf"],
    "tiff":  ["png", "jpg", "bmp",  "webp", "pdf"],
    "webp":  ["png", "jpg", "bmp",  "tiff", "pdf"],
    "gif":   ["png", "jpg", "bmp",  "pdf"],
    "pptx":  ["txt", "pdf"],
    "xml":   ["json", "txt"],
}

ALL_EXTS = sorted(FORMAT_TARGETS.keys())

ACCEPT_STRING = " ".join(f"*.{e}" for e in ALL_EXTS)


# ─── GUI ─────────────────────────────────────────────────────────────────────

class ConverterApp(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("File Converter")
        self.resizable(False, False)
        self.configure(bg="#f5f5f3")
        self._build_ui()
        self._src_path = None

    def _build_ui(self):
        PAD = {"padx": 20, "pady": 8}

        # Title
        tk.Label(self, text="File Converter", font=("Helvetica", 17, "bold"),
                 bg="#f5f5f3", fg="#1a1a1a").grid(row=0, column=0, columnspan=3,
                 padx=20, pady=(18, 2), sticky="w")
        tk.Label(self, text="Select a file, choose a target format, and convert.",
                 font=("Helvetica", 11), bg="#f5f5f3", fg="#666").grid(
                 row=1, column=0, columnspan=3, padx=20, pady=(0, 12), sticky="w")

        # File picker row
        tk.Label(self, text="Input file:", font=("Helvetica", 11),
                 bg="#f5f5f3", fg="#333").grid(row=2, column=0, sticky="w", padx=20)

        self.path_var = tk.StringVar(value="No file selected")
        self.path_label = tk.Label(self, textvariable=self.path_var,
                                   font=("Helvetica", 10), bg="#ffffff",
                                   fg="#444", relief="flat", bd=0,
                                   width=38, anchor="w",
                                   padx=8, pady=6,
                                   highlightbackground="#ccc",
                                   highlightthickness=1)
        self.path_label.grid(row=2, column=1, padx=(0, 8), pady=8)

        browse_btn = tk.Button(self, text="Browse…", command=self._browse,
                               font=("Helvetica", 10), bg="#ffffff",
                               fg="#222", relief="flat",
                               highlightbackground="#bbb",
                               highlightthickness=1,
                               padx=10, pady=5, cursor="hand2")
        browse_btn.grid(row=2, column=2, padx=(0, 20))

        # Format selection
        tk.Label(self, text="Convert to:", font=("Helvetica", 11),
                 bg="#f5f5f3", fg="#333").grid(row=3, column=0, sticky="w", padx=20, pady=4)

        self.fmt_var = tk.StringVar(value="")
        self.fmt_combo = ttk.Combobox(self, textvariable=self.fmt_var,
                                       state="disabled", width=20,
                                       font=("Helvetica", 11))
        self.fmt_combo.grid(row=3, column=1, sticky="w", padx=(0, 8), pady=4)

        # Progress bar
        self.progress = ttk.Progressbar(self, mode="indeterminate", length=340)
        self.progress.grid(row=4, column=0, columnspan=3, padx=20, pady=(10, 4))

        # Status label
        self.status_var = tk.StringVar(value="")
        self.status_lbl = tk.Label(self, textvariable=self.status_var,
                                    font=("Helvetica", 10), bg="#f5f5f3",
                                    fg="#555", wraplength=400, justify="left")
        self.status_lbl.grid(row=5, column=0, columnspan=3, padx=20, pady=(0, 4))

        # Convert button
        self.convert_btn = tk.Button(self, text="Convert",
                                      command=self._start_convert,
                                      font=("Helvetica", 12, "bold"),
                                      bg="#1a1a1a", fg="white",
                                      relief="flat", padx=20, pady=8,
                                      cursor="hand2", state="disabled")
        self.convert_btn.grid(row=6, column=0, columnspan=3, padx=20, pady=(8, 20))

        self.columnconfigure(1, weight=1)

    def _browse(self):
        path = filedialog.askopenfilename(
            title="Select a file to convert",
            filetypes=[
                ("Supported files", ACCEPT_STRING),
                ("All files", "*.*"),
            ]
        )
        if not path:
            return
        self._src_path = path
        short = os.path.basename(path)
        self.path_var.set(short if len(short) < 40 else "…" + short[-37:])

        ext = os.path.splitext(path)[1].lower().lstrip(".")
        targets = FORMAT_TARGETS.get(ext, [])

        if not targets:
            self._set_status(f"⚠ .{ext} files are not supported.", "orange")
            self.fmt_combo.configure(state="disabled")
            self.convert_btn.configure(state="disabled")
            return

        self.fmt_combo["values"] = [f".{t}  ({t.upper()})" for t in targets]
        self.fmt_combo.current(0)
        self.fmt_combo.configure(state="readonly")
        self.convert_btn.configure(state="normal")
        self._set_status(f"Ready — {len(targets)} output formats available.", "#555")

    def _start_convert(self):
        if not self._src_path or not self.fmt_var.get():
            return
        target = self.fmt_var.get().split()[0].lstrip(".")
        self.convert_btn.configure(state="disabled")
        self.fmt_combo.configure(state="disabled")
        self.progress.start(12)
        self._set_status("Converting…", "#555")
        thread = threading.Thread(target=self._run_convert, args=(self._src_path, target), daemon=True)
        thread.start()

    def _run_convert(self, src, target_fmt):
        out, err = convert_file(src, target_fmt)
        self.after(0, self._finish_convert, out, err)

    def _finish_convert(self, out, err):
        self.progress.stop()
        self.convert_btn.configure(state="normal")
        ext = os.path.splitext(self._src_path)[1].lower().lstrip(".")
        targets = FORMAT_TARGETS.get(ext, [])
        self.fmt_combo.configure(state="readonly")

        if err:
            self._set_status(f"❌  Error: {err}", "red")
            messagebox.showerror("Conversion failed", err)
        else:
            self._set_status(f"✅  Saved: {os.path.basename(out)}", "green")
            messagebox.showinfo("Done!",
                f"File converted successfully!\n\nSaved as:\n{out}")

    def _set_status(self, msg, color="#555"):
        self.status_var.set(msg)
        self.status_lbl.configure(fg=color)


# ─── Entry point ─────────────────────────────────────────────────────────────

if __name__ == "__main__":
    missing = check_deps()
    if missing:
        root = tk.Tk()
        root.withdraw()
        pkg_list = "\n  pip install " + " ".join(missing)
        messagebox.showwarning(
            "Missing dependencies",
            f"Some packages are not installed. Run:\n{pkg_list}\n\nThen restart the app."
        )
        root.destroy()

    app = ConverterApp()
    app.mainloop()